"""PPTX 파일 파싱 서비스 — 슬라이드별 텍스트, 이미지, 발표자 노트 추출."""

from __future__ import annotations

import base64
import logging
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.models.schemas import SlideContent

logger = logging.getLogger(__name__)


def parse_pptx(file_path: str | Path, output_dir: str | Path) -> list[SlideContent]:
    """PPTX 파일을 파싱하여 슬라이드별 콘텐츠를 추출한다.

    Parameters
    ----------
    file_path : 업로드된 .pptx 파일 경로
    output_dir : 추출된 이미지를 저장할 디렉토리

    Returns
    -------
    list[SlideContent]
    """
    file_path = Path(file_path)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(file_path))
    slides: list[SlideContent] = []

    for idx, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        image_paths: list[str] = []

        # 텍스트 & 이미지 추출 (애니메이션은 무시 — shape 정보만 사용)
        for shape in slide.shapes:
            # 텍스트 추출
            if shape.has_text_frame:
                full_text = "\n".join(
                    para.text for para in shape.text_frame.paragraphs if para.text.strip()
                )
                if full_text.strip():
                    texts.append(full_text)

            # 테이블 텍스트 추출
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            texts.append(cell_text)

            # 이미지 추출
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"
                img_filename = f"slide_{idx}_img_{len(image_paths) + 1}.{ext}"
                img_path = output_dir / img_filename
                img_path.write_bytes(image.blob)
                image_paths.append(str(img_path))

            # 그룹 shape 내부 처리
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                _extract_group(shape, texts, image_paths, output_dir, idx)

        # 발표자 노트 추출
        speaker_notes = ""
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            speaker_notes = notes_frame.text.strip()

        slides.append(
            SlideContent(
                slide_number=idx,
                texts=texts,
                speaker_notes=speaker_notes,
                image_paths=image_paths,
            )
        )

    logger.info("파싱 완료: %d개 슬라이드 추출", len(slides))
    return slides


def _extract_group(
    group_shape,
    texts: list[str],
    image_paths: list[str],
    output_dir: Path,
    slide_idx: int,
) -> None:
    """그룹 shape 내부의 텍스트와 이미지를 재귀적으로 추출한다."""
    for shape in group_shape.shapes:
        if shape.has_text_frame:
            full_text = "\n".join(
                para.text for para in shape.text_frame.paragraphs if para.text.strip()
            )
            if full_text.strip():
                texts.append(full_text)

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            ext = image.content_type.split("/")[-1]
            if ext == "jpeg":
                ext = "jpg"
            img_filename = f"slide_{slide_idx}_img_{len(image_paths) + 1}.{ext}"
            img_path = output_dir / img_filename
            img_path.write_bytes(image.blob)
            image_paths.append(str(img_path))

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _extract_group(shape, texts, image_paths, output_dir, slide_idx)


def encode_image_base64(image_path: str) -> str:
    """이미지 파일을 base64로 인코딩한다."""
    return base64.standard_b64encode(Path(image_path).read_bytes()).decode("utf-8")

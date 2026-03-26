"""테스트 공통 fixture — SQLite in-memory DB + mock 외부 API."""

from __future__ import annotations

import json
import os
import tempfile
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from sqlalchemy import create_engine, event
from sqlalchemy.orm import Session, sessionmaker
from sqlalchemy.pool import StaticPool

# pgvector 없는 SQLite에서도 모델이 로드되도록 환경 설정
os.environ.setdefault("DATABASE_URL", "sqlite:///:memory:")
os.environ.setdefault("ANTHROPIC_API_KEY", "test-key")
os.environ.setdefault("OPENAI_API_KEY", "test-key")
os.environ.setdefault("DEEPL_API_KEY", "test-key")

from app.database import Base, get_db  # noqa: E402
from app.main import app  # noqa: E402


# ---------------------------------------------------------------------------
# SQLite 용 engine (pgvector 컬럼은 TEXT로 폴백)
# StaticPool + check_same_thread=False → TestClient 스레드 문제 해결
# ---------------------------------------------------------------------------

TEST_ENGINE = create_engine(
    "sqlite:///:memory:",
    echo=False,
    connect_args={"check_same_thread": False},
    poolclass=StaticPool,
)
TestSession = sessionmaker(bind=TEST_ENGINE, autocommit=False, autoflush=False)


@event.listens_for(TEST_ENGINE, "connect")
def _set_sqlite_pragma(dbapi_conn, connection_record):
    cursor = dbapi_conn.cursor()
    cursor.execute("PRAGMA foreign_keys=ON")
    cursor.close()


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture(autouse=True)
def _create_tables():
    """각 테스트 전에 테이블을 생성하고 끝나면 삭제한다."""
    _patch_vector_column()
    Base.metadata.create_all(bind=TEST_ENGINE)
    yield
    Base.metadata.drop_all(bind=TEST_ENGINE)


@pytest.fixture()
def db():
    """테스트용 DB 세션."""
    session = TestSession()
    try:
        yield session
    finally:
        session.rollback()
        session.close()


@pytest.fixture()
def client(db):
    """FastAPI TestClient (DB 세션 오버라이드)."""
    from fastapi.testclient import TestClient

    def _override_get_db():
        try:
            yield db
        finally:
            pass

    app.dependency_overrides[get_db] = _override_get_db
    with TestClient(app) as c:
        yield c
    app.dependency_overrides.clear()


@pytest.fixture()
def tmp_dir():
    """임시 디렉토리."""
    with tempfile.TemporaryDirectory() as d:
        yield Path(d)


@pytest.fixture()
def sample_pptx(tmp_dir) -> Path:
    """텍스트 + 발표자 노트가 있는 샘플 PPTX 파일을 생성한다."""
    from pptx import Presentation

    prs = Presentation()
    for i in range(1, 4):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
        slide.shapes.title.text = f"슬라이드 {i} 제목"
        slide.placeholders[1].text = f"슬라이드 {i} 본문 내용입니다."
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = f"이것은 슬라이드 {i}의 발표자 노트입니다."

    path = tmp_dir / "test.pptx"
    prs.save(str(path))
    return path


# ---------------------------------------------------------------------------
# Mock helpers
# ---------------------------------------------------------------------------

def make_claude_response(text: str = "테스트 스크립트입니다."):
    """Anthropic Claude API 응답 mock 객체."""
    content_block = MagicMock()
    content_block.text = text
    usage = MagicMock()
    usage.input_tokens = 100
    usage.output_tokens = 50
    response = MagicMock()
    response.content = [content_block]
    response.usage = usage
    return response


def make_embedding_response(count: int = 1, dim: int = 8):
    """OpenAI 임베딩 응답 mock 객체 (작은 차원으로 테스트)."""
    items = []
    for i in range(count):
        item = MagicMock()
        item.embedding = [0.1 * (i + 1)] * dim
        items.append(item)
    response = MagicMock()
    response.data = items
    return response


# ---------------------------------------------------------------------------
# pgvector Vector 컬럼 → SQLite TEXT 패치
# ---------------------------------------------------------------------------

def _patch_vector_column():
    """pgvector의 Vector 타입을 SQLite 호환 TEXT로 교체한다."""
    from sqlalchemy import Text as SAText

    try:
        from app.models.embedding import SlideEmbedding
        for col in SlideEmbedding.__table__.columns:
            if col.name == "embedding":
                col.type = SAText()
    except Exception:
        pass

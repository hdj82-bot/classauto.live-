"""PPT 파싱 서비스 테스트."""

from pathlib import Path

from app.services.parser import encode_image_base64, parse_pptx


class TestParsePptx:
    """parse_pptx 함수 테스트."""

    def test_parse_extracts_correct_slide_count(self, sample_pptx, tmp_dir):
        output = tmp_dir / "output"
        slides = parse_pptx(sample_pptx, output)
        assert len(slides) == 3

    def test_parse_extracts_text_content(self, sample_pptx, tmp_dir):
        output = tmp_dir / "output"
        slides = parse_pptx(sample_pptx, output)
        # 각 슬라이드에 제목 + 본문 텍스트가 있어야 함
        for i, slide in enumerate(slides, start=1):
            assert slide.slide_number == i
            combined = " ".join(slide.texts)
            assert f"슬라이드 {i}" in combined

    def test_parse_extracts_speaker_notes(self, sample_pptx, tmp_dir):
        output = tmp_dir / "output"
        slides = parse_pptx(sample_pptx, output)
        for i, slide in enumerate(slides, start=1):
            assert f"슬라이드 {i}의 발표자 노트" in slide.speaker_notes

    def test_parse_creates_output_dir(self, sample_pptx, tmp_dir):
        output = tmp_dir / "new_output"
        assert not output.exists()
        parse_pptx(sample_pptx, output)
        assert output.exists()

    def test_parse_empty_slide_returns_empty_texts(self, tmp_dir):
        """빈 슬라이드만 있는 PPTX."""
        from pptx import Presentation

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        path = tmp_dir / "empty.pptx"
        prs.save(str(path))

        output = tmp_dir / "output"
        slides = parse_pptx(path, output)
        assert len(slides) == 1
        assert slides[0].texts == []
        assert slides[0].speaker_notes == ""


class TestEncodeImageBase64:
    """encode_image_base64 테스트."""

    def test_encodes_file_to_base64(self, tmp_dir):
        img = tmp_dir / "test.png"
        img.write_bytes(b"\x89PNG\r\n\x1a\n" + b"\x00" * 100)
        result = encode_image_base64(str(img))
        assert isinstance(result, str)
        assert len(result) > 0

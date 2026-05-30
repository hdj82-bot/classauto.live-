"""E2E 파이프라인 통합 테스트.

외부 서비스(S3, OpenAI, Claude, ElevenLabs, HeyGen)를 mock하고
전체 플로우를 검증합니다.

1. PPT 업로드 → S3 → Celery 5단계 파이프라인
2. TTS → S3 오디오 → HeyGen 렌더링
3. HeyGen 웹훅 → S3 비디오 → 완료
"""
import io
import uuid
from unittest.mock import patch, AsyncMock, MagicMock

import pytest
from pptx import Presentation

from app.services.pipeline.schemas import SlideScript
from tests.conftest import make_auth_header


# ── 헬퍼: 테스트용 PPTX 파일 생성 ───────────────────────────────────────────

def _create_test_pptx(slide_texts: list[str] | None = None) -> bytes:
    """실제 PPTX 바이트를 생성."""
    prs = Presentation()
    texts = slide_texts or ["안녕하세요, 첫 번째 슬라이드입니다.", "두 번째 슬라이드 내용입니다."]
    for text in texts:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ══════════════════════════════════════════════════════════════════════════════
# 1. PPT 업로드 → S3 → 파이프라인 시작 (API 레벨)
# ══════════════════════════════════════════════════════════════════════════════

@pytest.mark.asyncio
async def test_e2e_ppt_upload_triggers_pipeline(client, professor, lecture):
    """PPT 업로드 API가 S3 업로드 + Celery 파이프라인을 올바르게 트리거하는지 검증."""
    pptx_bytes = _create_test_pptx()

    with patch("app.services.pipeline.s3.upload_ppt", return_value=("https://s3/ppt/test.pptx", "ppt/lec/test.pptx")) as mock_s3, \
         patch("app.tasks.pipeline.start_pipeline") as mock_pipeline:
        mock_pipeline.return_value = MagicMock(id="celery-chain-001")

        resp = await client.post(
            "/api/v1/render/upload",
            params={"lecture_id": str(lecture.id)},
            files={"file": ("강의.pptx", pptx_bytes, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
            headers=make_auth_header(professor),
        )

    assert resp.status_code == 200
    data = resp.json()
    assert data["s3_url"] == "https://s3/ppt/test.pptx"
    assert data["celery_task_id"] == "celery-chain-001"

    # S3 업로드 검증
    mock_s3.assert_called_once()
    call_args = mock_s3.call_args
    assert call_args[0][1] == str(lecture.id)  # lecture_id
    # Critical 6: 사용자 제공 파일명("강의.pptx") 은 UUID.hex+.pptx 로 강제 치환됨
    safe_filename = call_args[0][2]
    assert safe_filename != "강의.pptx"
    assert safe_filename.endswith(".pptx")
    assert "/" not in safe_filename and ".." not in safe_filename
    assert len(safe_filename) == 32 + len(".pptx")  # uuid4().hex + ".pptx"

    # 파이프라인에 s3_key가 전달되는지 검증
    mock_pipeline.assert_called_once()
    pipeline_args = mock_pipeline.call_args[0]
    assert pipeline_args[1] == "ppt/lec/test.pptx"  # s3_key
    assert pipeline_args[2] == str(professor.id)  # instructor_id
    assert pipeline_args[3] == str(lecture.id)  # lecture_id


# ══════════════════════════════════════════════════════════════════════════════
# 2. Step1: S3 다운로드 → PPTX 파싱
# ══════════════════════════════════════════════════════════════════════════════

def test_e2e_step1_parse_from_s3(tmp_path):
    """step1_parse가 S3에서 다운로드 후 PPTX를 올바르게 파싱하는지 검증."""
    from app.tasks.pipeline import step1_parse

    pptx_bytes = _create_test_pptx(["슬라이드 A", "슬라이드 B", "슬라이드 C"])

    with patch("app.services.pipeline.s3.download_file", return_value=pptx_bytes) as mock_download, \
         patch("app.tasks.pipeline.UPLOAD_DIR", str(tmp_path)):
        # Celery task를 직접 호출 (task.run 대신 함수 직접 호출)
        result = step1_parse.run("test-task-001", "ppt/lecture-1/abc_test.pptx")

    mock_download.assert_called_once_with("ppt/lecture-1/abc_test.pptx")
    assert result["task_id"] == "test-task-001"
    assert result["s3_key"] == "ppt/lecture-1/abc_test.pptx"
    assert len(result["slides"]) == 3
    assert result["slides"][0]["slide_number"] == 1
    assert any("슬라이드 A" in t for t in result["slides"][0]["texts"])


# ══════════════════════════════════════════════════════════════════════════════
# 3. Step2: 임베딩 생성 (OpenAI mock)
# ══════════════════════════════════════════════════════════════════════════════

def test_e2e_step2_embed():
    """step2_embed가 슬라이드 텍스트를 임베딩으로 변환하는지 검증."""
    from app.tasks.pipeline import step2_embed

    prev_result = {
        "task_id": "test-task-001",
        "slides": [
            {"slide_number": 1, "texts": ["인공지능 개론"], "speaker_notes": "AI 소개", "image_paths": []},
            {"slide_number": 2, "texts": ["머신러닝 기초"], "speaker_notes": "", "image_paths": []},
        ],
    }

    fake_embeddings = [[0.1] * 1536, [0.2] * 1536]

    with patch("app.services.pipeline.embedding.get_embeddings", return_value=fake_embeddings), \
         patch("app.tasks.pipeline.SyncSessionLocal") as mock_session_cls:
        mock_db = MagicMock()
        mock_session_cls.return_value = mock_db

        result = step2_embed.run(prev_result)

    assert result["task_id"] == "test-task-001"
    assert len(result["slides"]) == 2
    mock_db.add_all.assert_called_once()
    mock_db.commit.assert_called_once()


# ══════════════════════════════════════════════════════════════════════════════
# 4. Step3: Claude 스크립트 생성
# ══════════════════════════════════════════════════════════════════════════════

def test_e2e_step3_generate_scripts():
    """step3가 Claude API로 스크립트를 생성하는지 검증."""
    from app.tasks.pipeline import step3_generate_scripts

    prev_result = {
        "task_id": "test-task-001",
        "slides": [
            {"slide_number": 1, "texts": ["파이썬 기초"], "speaker_notes": "파이썬 소개", "image_paths": []},
            {"slide_number": 2, "texts": ["변수와 타입"], "speaker_notes": "", "image_paths": []},
        ],
    }

    mock_scripts = [
        SlideScript(slide_number=1, script="안녕하세요, 오늘은 파이썬 기초를 배워보겠습니다."),
        SlideScript(slide_number=2, script="변수와 자료형에 대해 알아봅시다."),
    ]

    with patch("app.services.pipeline.script_generator.generate_scripts", return_value=mock_scripts):
        result = step3_generate_scripts.run(prev_result)

    assert "scripts" in result
    assert len(result["scripts"]) == 2
    assert result["scripts"][0]["script"] == "안녕하세요, 오늘은 파이썬 기초를 배워보겠습니다."
    assert result["scripts"][1]["slide_number"] == 2


# ══════════════════════════════════════════════════════════════════════════════
# 5. Step4 + Step5: 상태 마킹 + 알림
# ══════════════════════════════════════════════════════════════════════════════

def test_e2e_step4_persists_video_and_script():
    """step4 가 생성된 스크립트를 Video+VideoScript 로 영속화하고
    PENDING_REVIEW 로 마킹하는지 검증.

    종전 step4 는 DB 무작업이라 videos/video_scripts row 가 생성되지
    않아 studio 가 모든 강의에서 빈 스크립트(데모 폴백)를 보였다.
    """
    from app.tasks.pipeline import step4_mark_pending_review

    lecture_id = str(uuid.uuid4())
    prev_result = {
        "task_id": "test-task-001",
        "slides": [],
        "lecture_id": lecture_id,
        "scripts": [
            {"slide_number": 1, "script": "첫 번째 슬라이드 발화입니다."},
            {"slide_number": 2, "script": "두 번째 슬라이드 발화입니다."},
        ],
    }

    with patch("app.tasks.pipeline.SyncSessionLocal") as mock_session_cls:
        mock_db = MagicMock()
        mock_session_cls.return_value = mock_db
        # Video / VideoScript 둘 다 미존재 → 생성 경로
        mock_db.execute.return_value.scalars.return_value.first.return_value = None

        result = step4_mark_pending_review.run(prev_result)

    assert result["status"] == "PENDING_REVIEW"
    assert "video_id" in result
    # Video + VideoScript 두 row add, 단일 트랜잭션 commit
    assert mock_db.add.call_count == 2
    mock_db.commit.assert_called_once()


def test_e2e_step4_missing_lecture_id_raises():
    """lecture_id 누락 시 조용히 통과하지 않고 실패해야 한다
    (종전엔 DB 무작업이라 누락이어도 통과 — 영속화 누락 은폐)."""
    from app.tasks.pipeline import step4_mark_pending_review

    with pytest.raises(RuntimeError, match="lecture_id"):
        step4_mark_pending_review.run(
            {"task_id": "t", "slides": [], "scripts": []}
        )


def test_estimate_segments_timing_and_shape():
    """_estimate_segments: 텍스트 길이 기반 누적 타이밍 + 0-based
    slide_index + 필수 필드(ScriptSegment 호환) 검증."""
    from app.tasks.pipeline import _estimate_segments

    scripts = [
        {"slide_number": 2, "script": "가" * 50},   # 50/5=10초
        {"slide_number": 1, "script": "나" * 3},    # min 5초
    ]
    segs = _estimate_segments(scripts)

    # slide_number 정렬 → slide_index 0-based
    assert [s["slide_index"] for s in segs] == [0, 1]
    # 슬라이드1: 0~5(min), 슬라이드2: 5~15(누적)
    assert (segs[0]["start_seconds"], segs[0]["end_seconds"]) == (0, 5)
    assert (segs[1]["start_seconds"], segs[1]["end_seconds"]) == (5, 15)
    for s in segs:
        assert s["tone"] == "normal"
        assert s["question_pin_seconds"] is None
        assert s["text"]


def test_e2e_step5_notify():
    from app.tasks.pipeline import step5_notify

    instructor_id = str(uuid.uuid4())
    lecture_id = str(uuid.uuid4())
    prev_result = {
        "task_id": "test-task-001",
        "instructor_id": instructor_id,
        "lecture_id": lecture_id,
        "status": "PENDING_REVIEW",
    }

    with patch("app.services.pipeline.notification.notify_instructor", new_callable=AsyncMock) as mock_notify:
        result = step5_notify.run(prev_result)

    assert result["task_id"] == "test-task-001"
    mock_notify.assert_called_once()


# ══════════════════════════════════════════════════════════════════════════════
# 6. 렌더링 태스크: TTS → S3 → HeyGen
# ══════════════════════════════════════════════════════════════════════════════

def test_e2e_render_slide_pipeline():
    """render_slide 태스크가 TTS → S3 업로드 → HeyGen 요청 전체 플로우를 실행하는지 검증."""
    from app.tasks.render import render_slide
    from app.models.video_render import RenderStatus
    from app.services.pipeline.tts import TTSResult

    render_id = uuid.uuid4()
    instructor_id = uuid.uuid4()
    mock_render = MagicMock()
    mock_render.id = render_id
    mock_render.instructor_id = instructor_id
    mock_render.avatar_id = "avatar-test"
    mock_render.status = RenderStatus.pending
    # Critical 8: idempotent skip 분기 비활성화 — 신규 렌더 상태로 진입
    mock_render.audio_url = None
    mock_render.heygen_job_id = None

    # render task 가 lecture 에서 voice_gender / voice_id / voice_speed / avatar_scale 를
    # lookup 한다. mock 으로 명시하지 않으면 MagicMock 객체가 그대로 흘러 synthesize/
    # create_video 인자에 들어가 assert 가 깨지므로 구체값으로 고정한다.
    mock_lecture = MagicMock()
    mock_lecture.voice_gender = "male"
    mock_lecture.voice_id = None
    mock_lecture.voice_speed = 1.0
    mock_lecture.avatar_scale = 1.0

    mock_tts_result = TTSResult(audio_bytes=b"tts-audio", provider="elevenlabs", duration_seconds=1.2)

    with patch("app.tasks.render.SyncSessionLocal") as mock_session_cls, \
         patch("app.services.pipeline.tts.synthesize", new_callable=AsyncMock, return_value=mock_tts_result) as mock_tts, \
         patch("app.services.pipeline.s3.upload_audio_bytes", return_value="https://s3/audio/test.mp3") as mock_s3_audio, \
         patch("app.services.pipeline.s3.file_exists", return_value=False), \
         patch("app.services.pipeline.heygen.create_video", new_callable=AsyncMock, return_value="heygen-vid-001") as mock_heygen, \
         patch("app.services.pipeline.cost_log.record_once"), \
         patch("app.services.pipeline.cost_log.record"):
        mock_db = MagicMock()
        mock_db.query.return_value.filter.return_value.one.return_value = mock_render
        # render task 의 두 번째 query — Lecture lookup (first()) 도 명시 매핑.
        mock_db.query.return_value.filter.return_value.first.return_value = mock_lecture
        mock_session_cls.return_value = mock_db

        # Critical 7: caller_user_id 일치 시 통과 — instructor_id 그대로 전달
        result = render_slide.apply(
            args=[str(render_id), "안녕하세요, 테스트 스크립트입니다.", str(instructor_id)],
        ).get(propagate=True)

    # TTS 호출 검증 — 교수자가 고른 voice_id·voice_speed 와 gender 가 함께 전달된다.
    mock_tts.assert_called_once_with(
        "안녕하세요, 테스트 스크립트입니다.",
        voice_id=None,
        gender="male",
        speed=1.0,
        cloned=False,
    )

    # S3 오디오 업로드 검증
    mock_s3_audio.assert_called_once_with(b"tts-audio", str(render_id))

    # HeyGen 비디오 생성 검증 — 0016 이후 gender, #239 이후 avatar_scale 함께 전달
    mock_heygen.assert_called_once_with(
        audio_url="https://s3/audio/test.mp3",
        avatar_id="avatar-test",
        gender="male",
        callback_id=str(render_id),
        avatar_scale=1.0,
    )

    # 결과 검증
    assert result["render_id"] == str(render_id)
    assert result["heygen_job_id"] == "heygen-vid-001"
    assert mock_render.heygen_job_id == "heygen-vid-001"
    assert mock_render.audio_url == "https://s3/audio/test.mp3"


# ══════════════════════════════════════════════════════════════════════════════
# 7. 전체 API 플로우: 업로드 → 렌더링 요청 → 상태 조회
# ══════════════════════════════════════════════════════════════════════════════

@pytest.mark.asyncio
async def test_e2e_full_api_flow(client, professor, lecture, db):
    """업로드 → 렌더링 요청 → 상태 조회 API 전체 플로우."""

    # Step 1: PPT 업로드
    pptx_bytes = _create_test_pptx()
    with patch("app.services.pipeline.s3.upload_ppt", return_value=("https://s3/test.pptx", "ppt/key")), \
         patch("app.tasks.pipeline.start_pipeline") as mock_pipeline:
        mock_pipeline.return_value = MagicMock(id="chain-001")
        upload_resp = await client.post(
            "/api/v1/render/upload",
            params={"lecture_id": str(lecture.id)},
            files={"file": ("test.pptx", pptx_bytes, "application/octet-stream")},
            headers=make_auth_header(professor),
        )
    assert upload_resp.status_code == 200

    # Step 2: 렌더링 요청
    with patch("app.services.pipeline.subscription.check_limit", new_callable=AsyncMock), \
         patch("app.tasks.render.render_slide") as mock_task:
        mock_task.delay = MagicMock()
        render_resp = await client.post(
            "/api/v1/render",
            params={"lecture_id": str(lecture.id)},
            json=[
                {"script": "첫 번째 슬라이드 스크립트", "slide_number": 1},
                {"script": "두 번째 슬라이드 스크립트", "slide_number": 2},
            ],
            headers=make_auth_header(professor),
        )
    assert render_resp.status_code == 200
    render_data = render_resp.json()
    assert len(render_data["render_ids"]) == 2

    # Step 3: 렌더 상태 조회
    status_resp = await client.get(
        f"/api/v1/render/lecture/{lecture.id}",
        headers=make_auth_header(professor),
    )
    assert status_resp.status_code == 200
    status_data = status_resp.json()
    assert status_data["total"] == 2
    assert status_data["lecture_id"] == str(lecture.id)


# ══════════════════════════════════════════════════════════════════════════════
# 8. PPTX 파서 단위 테스트 (실제 파일)
# ══════════════════════════════════════════════════════════════════════════════

def test_parser_real_pptx(tmp_path):
    """실제 PPTX를 생성하고 파서가 올바르게 추출하는지 검증."""
    from app.services.pipeline.parser import parse_pptx

    pptx_bytes = _create_test_pptx(["데이터 사이언스 개론", "통계학 기초", "머신러닝 알고리즘"])
    pptx_path = tmp_path / "test.pptx"
    pptx_path.write_bytes(pptx_bytes)
    output_dir = tmp_path / "images"

    slides = parse_pptx(str(pptx_path), str(output_dir))

    assert len(slides) == 3
    assert slides[0].slide_number == 1
    assert any("데이터 사이언스" in t for t in slides[0].texts)
    assert slides[2].slide_number == 3
    assert any("머신러닝" in t for t in slides[2].texts)


# ══════════════════════════════════════════════════════════════════════════════
# 9. 폴링 태스크 E2E
# ══════════════════════════════════════════════════════════════════════════════

def test_e2e_polling_completed_render():
    """poll_pending_renders가 완료된 렌더를 S3 업로드 + READY 처리하는지 검증."""
    from app.models.video_render import RenderStatus
    from app.tasks.polling import poll_pending_renders

    render_id = uuid.uuid4()
    lecture_id = uuid.uuid4()
    instructor_id = uuid.uuid4()

    mock_render = MagicMock()
    mock_render.id = render_id
    mock_render.heygen_job_id = "heygen-job-done"
    mock_render.lecture_id = lecture_id
    mock_render.instructor_id = instructor_id
    mock_render.slide_number = 1
    mock_render.status = RenderStatus.rendering

    mock_status = {
        "status": "completed",
        "video_url": "https://heygen.com/video.mp4",
        "duration": 25.0,
    }

    with patch("app.tasks.polling.SyncSessionLocal") as mock_session_cls, \
         patch("app.tasks.polling.get_video_status", new_callable=AsyncMock, return_value=mock_status), \
         patch("app.tasks.polling.s3_svc.upload_from_url", new_callable=AsyncMock, return_value=("https://s3/video.mp4", 3.5)), \
         patch("app.tasks.polling.cost_log.record_once"), \
         patch("app.tasks.polling.notification.notify_instructor", new_callable=AsyncMock) as mock_notify:
        mock_db = MagicMock()
        mock_db.query.return_value.filter.return_value.all.return_value = [mock_render]
        mock_session_cls.return_value = mock_db

        result = poll_pending_renders.run()

    assert result["checked"] == 1
    assert result["completed"] == 1
    assert mock_render.s3_video_url == "https://s3/video.mp4"
    assert mock_render.status == RenderStatus.ready
    mock_notify.assert_called_once()

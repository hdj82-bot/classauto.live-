"""step1_parse 가 슬라이드 PNG 를 렌더·업로드하고 step2_embed 가 그 URL 을
SlideEmbedding 에 기록하는지 검증.

본 PR(창 2) 은 ``app.services.pipeline.slide_renderer`` 모듈과 ``s3.upload_slide_image``
헬퍼를 직접 만들지 않는다 — 창 1 PR 이 제공한다. 단, 단독 실행 가능하도록
fixture 가 동일 시그니처의 stub 을 sys.modules / s3 모듈에 주입한다.

[검증 항목]
- step1_parse 가 render_pptx_to_images + upload_slide_image 를 호출하고 결과를
  ``prev_result["slide_image_urls"]`` 로 내려준다.
- render 단계 실패 (SlideRenderError) 는 graceful — 파이프라인은 계속 진행.
- upload 단계 한 장 실패는 graceful — 나머지 슬라이드는 정상 처리.
- lecture_id 가 없으면 렌더 자체를 skip.
- step2_embed 가 ``slide_image_urls`` 를 ``store_slide_embeddings`` 에 전달.
"""
from __future__ import annotations

import io
import sys
import types
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from pptx import Presentation


# ── 헬퍼 ─────────────────────────────────────────────────────────────────────


def _make_pptx_bytes(slide_titles: list[str]) -> bytes:
    """간단한 PPTX 바이트 생성 — 제목만 들어간 슬라이드 N장."""
    prs = Presentation()
    for title in slide_titles:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Fixtures: 창 1 의 모듈/헬퍼 stub ─────────────────────────────────────────


@pytest.fixture
def stub_slide_renderer(monkeypatch):
    """``app.services.pipeline.slide_renderer`` 를 stub 모듈로 대체.

    실제 LibreOffice 호출 없이 ``input pptx 의 슬라이드 수와 정확히 같은`` PNG
    경로 리스트를 반환하도록 흉내낸다. ``SlideRenderError`` 도 같이 노출해
    step1 의 except 블록이 정확히 그 타입을 잡는 것을 검증할 수 있게 한다.
    """
    stub = types.ModuleType("app.services.pipeline.slide_renderer")

    class SlideRenderError(RuntimeError):
        pass

    stub.SlideRenderError = SlideRenderError

    call_log: list[dict] = []

    def fake_render(pptx_path, output_dir, dpi=110):
        out = Path(output_dir)
        out.mkdir(parents=True, exist_ok=True)
        # 페이지 수는 caller 가 stub 의 attribute 로 override 가능. 기본 3장.
        n = getattr(stub, "_n_pages", 3)
        paths: list[Path] = []
        for i in range(1, n + 1):
            p = out / f"slide-{i}.png"
            p.write_bytes(b"\x89PNG\r\n\x1a\n" + bytes([i] * 16))
            paths.append(p)
        call_log.append({"pptx_path": str(pptx_path), "n": n, "dpi": dpi})
        return paths

    stub.render_pptx_to_images = fake_render
    stub.call_log = call_log
    monkeypatch.setitem(sys.modules, "app.services.pipeline.slide_renderer", stub)
    yield stub


@pytest.fixture
def stub_upload_slide_image(monkeypatch):
    """``s3.upload_slide_image`` 를 빈 헬퍼로 대체. (lecture_id, slide_number, len) 기록."""
    from app.services.pipeline import s3 as s3_mod

    uploads: list[tuple[str, int, int]] = []

    def fake_upload(image_bytes: bytes, lecture_id: str, slide_number: int) -> str:
        uploads.append((lecture_id, slide_number, len(image_bytes)))
        return f"https://test-bucket.s3.amazonaws.com/slides/{lecture_id}/{slide_number}.png"

    monkeypatch.setattr(s3_mod, "upload_slide_image", fake_upload, raising=False)
    return uploads


# ── step1_parse 통합 ────────────────────────────────────────────────────────


def test_step1_renders_and_uploads_slide_images(
    tmp_path, stub_slide_renderer, stub_upload_slide_image
):
    """lecture_id 가 있으면 render + upload 가 호출되고 prev_result 에 URL dict 가 담긴다."""
    from app.tasks.pipeline import step1_parse

    pptx_bytes = _make_pptx_bytes(["A", "B", "C"])
    stub_slide_renderer._n_pages = 3

    with patch("app.services.pipeline.s3.download_file", return_value=pptx_bytes), \
            patch("app.tasks.pipeline.UPLOAD_DIR", str(tmp_path)):
        result = step1_parse.run(
            "task-img-1",
            "ppt/lec-x/file.pptx",
            instructor_id="prof-1",
            lecture_id="lec-x",
        )

    # prev_result 의 새 키
    assert "slide_image_urls" in result
    assert result["slide_image_urls"] == {
        1: "https://test-bucket.s3.amazonaws.com/slides/lec-x/1.png",
        2: "https://test-bucket.s3.amazonaws.com/slides/lec-x/2.png",
        3: "https://test-bucket.s3.amazonaws.com/slides/lec-x/3.png",
    }
    # render 가 정확히 한 번 호출
    assert len(stub_slide_renderer.call_log) == 1
    # 3장 모두 upload 호출
    assert {u[1] for u in stub_upload_slide_image} == {1, 2, 3}
    assert {u[0] for u in stub_upload_slide_image} == {"lec-x"}
    # 기존 키도 유지
    assert result["task_id"] == "task-img-1"
    assert result["lecture_id"] == "lec-x"
    assert len(result["slides"]) == 3


def test_step1_render_failure_is_graceful(
    tmp_path, stub_slide_renderer, stub_upload_slide_image
):
    """SlideRenderError 발생 시 파이프라인은 계속 진행, slide_image_urls 는 빈 dict."""
    from app.tasks.pipeline import step1_parse

    def boom(pptx_path, output_dir, dpi=110):
        raise stub_slide_renderer.SlideRenderError("soffice 가 죽었어요")

    stub_slide_renderer.render_pptx_to_images = boom

    pptx_bytes = _make_pptx_bytes(["X", "Y"])
    with patch("app.services.pipeline.s3.download_file", return_value=pptx_bytes), \
            patch("app.tasks.pipeline.UPLOAD_DIR", str(tmp_path)):
        result = step1_parse.run(
            "task-render-fail",
            "ppt/lec/render-fail.pptx",
            lecture_id="lec-render-fail",
        )

    assert result["slide_image_urls"] == {}
    # 파이프라인 자체는 정상 — 파싱 결과는 살아 있다
    assert len(result["slides"]) == 2
    # upload 는 한 번도 호출되지 않았어야 한다
    assert stub_upload_slide_image == []


def test_step1_upload_failure_skips_only_that_slide(
    tmp_path, stub_slide_renderer, monkeypatch
):
    """한 장의 upload 만 실패하면 그 슬라이드만 url dict 에서 빠진다 — 다른 슬라이드는 그대로."""
    from app.services.pipeline import s3 as s3_mod
    from app.tasks.pipeline import step1_parse

    stub_slide_renderer._n_pages = 3

    def selective_upload(image_bytes, lecture_id, slide_number):
        if slide_number == 2:
            raise RuntimeError("S3 5xx")
        return f"https://test/slides/{lecture_id}/{slide_number}.png"

    monkeypatch.setattr(s3_mod, "upload_slide_image", selective_upload, raising=False)

    pptx_bytes = _make_pptx_bytes(["A", "B", "C"])
    with patch("app.services.pipeline.s3.download_file", return_value=pptx_bytes), \
            patch("app.tasks.pipeline.UPLOAD_DIR", str(tmp_path)):
        result = step1_parse.run(
            "task-partial-upload",
            "ppt/lec/partial.pptx",
            lecture_id="lec-partial",
        )

    assert result["slide_image_urls"] == {
        1: "https://test/slides/lec-partial/1.png",
        3: "https://test/slides/lec-partial/3.png",
    }


def test_step1_no_lecture_id_skips_render(
    tmp_path, stub_slide_renderer, stub_upload_slide_image
):
    """lecture_id 가 없으면 (e.g. 기존 호출 경로) render 자체를 skip 한다."""
    from app.tasks.pipeline import step1_parse

    pptx_bytes = _make_pptx_bytes(["X"])
    with patch("app.services.pipeline.s3.download_file", return_value=pptx_bytes), \
            patch("app.tasks.pipeline.UPLOAD_DIR", str(tmp_path)):
        result = step1_parse.run("task-no-lec", "ppt/lec/none.pptx")

    assert result["slide_image_urls"] == {}
    assert stub_slide_renderer.call_log == []
    assert stub_upload_slide_image == []


# ── step2_embed 가 slide_image_urls 를 store_slide_embeddings 로 전달 ──────


def test_step2_passes_slide_image_urls_to_store():
    """prev_result.slide_image_urls 가 store_slide_embeddings 의 kwarg 로 전달된다."""
    from app.tasks.pipeline import step2_embed

    prev_result = {
        "task_id": "task-img-2",
        "slides": [
            {"slide_number": 1, "texts": ["A"], "speaker_notes": "", "image_paths": []},
            {"slide_number": 2, "texts": ["B"], "speaker_notes": "", "image_paths": []},
        ],
        "slide_image_urls": {
            1: "https://test/slides/L/1.png",
            2: "https://test/slides/L/2.png",
        },
    }

    with patch("app.services.pipeline.embedding.store_slide_embeddings") as mock_store, \
            patch("app.tasks.pipeline.SyncSessionLocal") as mock_session_cls:
        mock_store.return_value = 2
        mock_db = MagicMock()
        mock_session_cls.return_value = mock_db

        result = step2_embed.run(prev_result)

    # store_slide_embeddings 호출 시 slide_image_urls 가 kwarg 로 전달됐는지
    mock_store.assert_called_once()
    _args, kwargs = mock_store.call_args
    assert kwargs.get("slide_image_urls") == {
        1: "https://test/slides/L/1.png",
        2: "https://test/slides/L/2.png",
    }
    # prev_result 통과 (다음 단계로 흘러감)
    assert result["task_id"] == "task-img-2"


def test_step2_handles_missing_slide_image_urls_key():
    """slide_image_urls 키가 prev_result 에 아예 없어도 빈 dict 로 graceful 진행."""
    from app.tasks.pipeline import step2_embed

    prev_result = {
        "task_id": "task-missing-urls",
        "slides": [
            {"slide_number": 1, "texts": ["A"], "speaker_notes": "", "image_paths": []},
        ],
    }

    with patch("app.services.pipeline.embedding.store_slide_embeddings") as mock_store, \
            patch("app.tasks.pipeline.SyncSessionLocal") as mock_session_cls:
        mock_store.return_value = 1
        mock_session_cls.return_value = MagicMock()

        step2_embed.run(prev_result)

    _args, kwargs = mock_store.call_args
    assert kwargs.get("slide_image_urls") == {}


def test_step2_normalizes_str_keys_to_int():
    """Celery JSON 직렬화로 dict key 가 str 이 된 경우에도 int 로 정규화돼야 한다."""
    from app.tasks.pipeline import step2_embed

    prev_result = {
        "task_id": "task-str-keys",
        "slides": [
            {"slide_number": 1, "texts": ["A"], "speaker_notes": "", "image_paths": []},
        ],
        # JSON 직렬화 후 dict key 는 string 이 된다
        "slide_image_urls": {"1": "https://test/slides/L/1.png"},
    }

    with patch("app.services.pipeline.embedding.store_slide_embeddings") as mock_store, \
            patch("app.tasks.pipeline.SyncSessionLocal") as mock_session_cls:
        mock_store.return_value = 1
        mock_session_cls.return_value = MagicMock()

        step2_embed.run(prev_result)

    _args, kwargs = mock_store.call_args
    assert kwargs.get("slide_image_urls") == {1: "https://test/slides/L/1.png"}


# ── store_slide_embeddings 단위 ─────────────────────────────────────────────


def test_store_slide_embeddings_sets_slide_image_url_on_record():
    """slide_image_urls dict 의 url 이 SlideEmbedding.slide_image_url 에 매핑된다."""
    from app.services.pipeline.embedding import store_slide_embeddings
    from app.services.pipeline.schemas import SlideContent

    db = MagicMock()
    fake_emb = [0.1] * 1536

    with patch(
        "app.services.pipeline.embedding.get_embeddings",
        return_value=[fake_emb, fake_emb],
    ):
        slides = [
            SlideContent(slide_number=1, texts=["A"], speaker_notes=""),
            SlideContent(slide_number=2, texts=["B"], speaker_notes=""),
        ]
        store_slide_embeddings(
            db,
            "task-set-url",
            slides,
            slide_image_urls={1: "https://test/1.png"},
        )

    records = db.add_all.call_args[0][0]
    assert len(records) == 2
    by_slide = {r.slide_number: r for r in records}
    assert by_slide[1].slide_image_url == "https://test/1.png"
    # 누락된 슬라이드는 None
    assert by_slide[2].slide_image_url is None


def test_store_slide_embeddings_without_image_urls_defaults_to_none():
    """slide_image_urls kwarg 를 안 넘기면 모든 row 의 slide_image_url 이 None."""
    from app.services.pipeline.embedding import store_slide_embeddings
    from app.services.pipeline.schemas import SlideContent

    db = MagicMock()
    with patch(
        "app.services.pipeline.embedding.get_embeddings",
        return_value=[[0.1] * 1536],
    ):
        store_slide_embeddings(
            db,
            "task-no-urls",
            [SlideContent(slide_number=1, texts=["A"], speaker_notes="")],
        )

    records = db.add_all.call_args[0][0]
    assert records[0].slide_image_url is None

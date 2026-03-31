"""PPTX 파일 파싱 서비스."""
from __future__ import annotations

import base64
import logging
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.pipeline.schemas import SlideContent

logger = logging.getLogger(__name__)


def parse_pptx(file_path: str | Path, output_dir: str | Path) -> list[SlideContent]:
    """PPTX 파일을 파싱하여 슬라이드별 콘텐츠를 추출."""
    file_path = Path(file_path)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(file_path))
    slides: list[SlideContent] = []

    for idx, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        image_paths: list[str] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                full_text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
                if full_text.strip():
                    texts.append(full_text)

            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            texts.append(cell_text)

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"
                img_filename = f"slide_{idx}_img_{len(image_paths) + 1}.{ext}"
                img_path = output_dir / img_filename
                img_path.write_bytes(image.blob)
                image_paths.append(str(img_path))

            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                _extract_group(shape, texts, image_paths, output_dir, idx)

        speaker_notes = ""
        if slide.has_notes_slide:
            speaker_notes = slide.notes_slide.notes_text_frame.text.strip()

        slides.append(SlideContent(slide_number=idx, texts=texts, speaker_notes=speaker_notes, image_paths=image_paths))

    logger.info("파싱 완료: %d개 슬라이드 추출", len(slides))
    return slides


def _extract_group(group_shape, texts: list[str], image_paths: list[str], output_dir: Path, slide_idx: int) -> None:
    for shape in group_shape.shapes:
        if shape.has_text_frame:
            full_text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
            if full_text.strip():
                texts.append(full_text)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            ext = image.content_type.split("/")[-1]
            if ext == "jpeg":
                ext = "jpg"
            img_filename = f"slide_{slide_idx}_img_{len(image_paths) + 1}.{ext}"
            img_path = output_dir / img_filename
            img_path.write_bytes(image.blob)
            image_paths.append(str(img_path))
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _extract_group(shape, texts, image_paths, output_dir, slide_idx)


def encode_image_base64(image_path: str) -> str:
    return base64.standard_b64encode(Path(image_path).read_bytes()).decode("utf-8")
